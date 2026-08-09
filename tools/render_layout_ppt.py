#!/usr/bin/env python3
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR


SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
DEFAULT_CANVAS_W = 1600
DEFAULT_CANVAS_H = 900


def rgb(hex_color, default="333333"):
    s = (hex_color or default).strip().replace("#", "")
    if len(s) != 6:
        s = default
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def align_value(v):
    v = (v or "left").lower()
    if v == "center":
        return PP_ALIGN.CENTER
    if v == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def add_text_to_shape(shape, text, font_size=14, color="#333333", bold=False, font_name="Microsoft YaHei", align="left"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    lines = str(text or "").split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align_value(align)
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bool(bold)
            run.font.color.rgb = rgb(color)
            run.font.name = font_name


def render_ppt(layout_json, output_pptx):
    layout_path = Path(layout_json)
    out_path = Path(output_pptx)
    data = json.loads(layout_path.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    blank = prs.slide_layouts[6]
    slides = data.get("slides") or [data]

    for slide_data in slides:
        slide = prs.slides.add_slide(blank)

        cw = float(slide_data.get("canvas_width", DEFAULT_CANVAS_W))
        ch = float(slide_data.get("canvas_height", DEFAULT_CANVAS_H))

        def X(v): return Inches(float(v) / cw * SLIDE_W_IN)
        def Y(v): return Inches(float(v) / ch * SLIDE_H_IN)
        def W(v): return Inches(float(v) / cw * SLIDE_W_IN)
        def H(v): return Inches(float(v) / ch * SLIDE_H_IN)

        bg = slide_data.get("background", "#FFFFFF")
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(bg, "FFFFFF")

        for el in slide_data.get("elements", []):
            typ = el.get("type", "text")
            x = el.get("x", 0)
            y = el.get("y", 0)
            w = el.get("w", 100)
            h = el.get("h", 40)

            if typ in ["rect", "round_rect", "lane", "footer_bar"]:
                shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if typ == "round_rect" or el.get("rounded") else MSO_SHAPE.RECTANGLE
                shape = slide.shapes.add_shape(shape_type, X(x), Y(y), W(w), H(h))

                fill = el.get("fill", "#FFFFFF")
                if fill == "none":
                    shape.fill.background()
                else:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = rgb(fill, "FFFFFF")

                line_color = el.get("line", el.get("stroke", "#DDDDDD"))
                shape.line.color.rgb = rgb(line_color, "DDDDDD")
                shape.line.width = Pt(float(el.get("line_width", 1)))

                if el.get("text"):
                    add_text_to_shape(
                        shape,
                        el.get("text", ""),
                        font_size=el.get("font_size", 14),
                        color=el.get("color", "#333333"),
                        bold=el.get("bold", False),
                        font_name=el.get("font", "Microsoft YaHei"),
                        align=el.get("align", "center"),
                    )

            elif typ == "text":
                shape = slide.shapes.add_textbox(X(x), Y(y), W(w), H(h))
                add_text_to_shape(
                    shape,
                    el.get("text", ""),
                    font_size=el.get("font_size", 14),
                    color=el.get("color", "#333333"),
                    bold=el.get("bold", False),
                    font_name=el.get("font", "Microsoft YaHei"),
                    align=el.get("align", "left"),
                )

            elif typ == "card":
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, X(x), Y(y), W(w), H(h))
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb(el.get("fill", "#FFFFFF"), "FFFFFF")
                shape.line.color.rgb = rgb(el.get("line", "#DDDDDD"), "DDDDDD")
                shape.line.width = Pt(float(el.get("line_width", 1)))

                title = el.get("title", "")
                body = el.get("body", "")
                bullets = el.get("bullets", [])
                text = title
                if body:
                    text += "\n" + body
                if bullets:
                    text += "\n" + "\n".join([f"• {b}" for b in bullets])

                add_text_to_shape(
                    shape,
                    text,
                    font_size=el.get("font_size", 12),
                    color=el.get("color", "#333333"),
                    bold=el.get("bold", False),
                    font_name=el.get("font", "Microsoft YaHei"),
                    align=el.get("align", "left"),
                )

            elif typ == "line":
                x1 = el.get("x1", x)
                y1 = el.get("y1", y)
                x2 = el.get("x2", x + w)
                y2 = el.get("y2", y + h)
                conn = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    X(x1), Y(y1), X(x2), Y(y2)
                )
                conn.line.color.rgb = rgb(el.get("color", "#999999"), "999999")
                conn.line.width = Pt(float(el.get("width", 1)))

            elif typ == "image":
                path = el.get("path")
                if path and Path(path).exists():
                    slide.shapes.add_picture(path, X(x), Y(y), W(w), H(h))

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(str(out_path))


def main():
    if len(sys.argv) != 3:
        print("Usage: render_layout_ppt.py layout.json output.pptx", file=sys.stderr)
        sys.exit(2)

    render_ppt(sys.argv[1], sys.argv[2])


if __name__ == "__main__":
    main()
