from pathlib import Path
from typing import Dict, List, Tuple

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

ORANGE = RGBColor(240, 90, 35)
DARK = RGBColor(51, 51, 51)
GRAY = RGBColor(102, 102, 102)
WHITE = RGBColor(255, 255, 255)


def _fit_contain(box_left, box_top, box_w, box_h, img_w, img_h):
    box_ratio = box_w / box_h
    img_ratio = img_w / img_h

    if img_ratio > box_ratio:
        w = box_w
        h = int(w / img_ratio)
        x = box_left
        y = int(box_top + (box_h - h) / 2)
    else:
        h = box_h
        w = int(h * img_ratio)
        x = int(box_left + (box_w - w) / 2)
        y = box_top

    return int(x), int(y), int(w), int(h)


def _style_of(kind: str):
    kind = (kind or "").lower()
    if kind == "title":
        return {"size": 28, "bold": True, "color": DARK, "fill": WHITE, "transparency": 0.18, "align": PP_ALIGN.LEFT}
    if kind == "subtitle":
        return {"size": 13, "bold": False, "color": GRAY, "fill": WHITE, "transparency": 0.18, "align": PP_ALIGN.LEFT}
    if kind == "section":
        return {"size": 20, "bold": True, "color": ORANGE, "fill": WHITE, "transparency": 0.10, "align": PP_ALIGN.CENTER}
    if kind == "card_title":
        return {"size": 16, "bold": True, "color": ORANGE, "fill": WHITE, "transparency": 0.10, "align": PP_ALIGN.LEFT}
    if kind == "footer":
        return {"size": 18, "bold": True, "color": WHITE, "fill": ORANGE, "transparency": 0.0, "align": PP_ALIGN.CENTER}
    return {"size": 14, "bold": False, "color": DARK, "fill": WHITE, "transparency": 0.12, "align": PP_ALIGN.LEFT}


def _add_overlay(slide, x, y, w, h, text: str, kind: str):
    if not text:
        return

    st = _style_of(kind)

    # 背景遮罩
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        int(x), int(y), int(w), int(h)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = st["fill"]
    try:
        fill.transparency = st["transparency"]
    except Exception:
        pass
    shape.line.fill.background()

    # 文本
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = st["align"]
    run = p.add_run()
    run.text = text
    run.font.name = "楷体"
    run.font.size = Pt(st["size"])
    run.font.bold = st["bold"]
    run.font.color.rgb = st["color"]


def _default_overlays(layout: Dict):
    overlays = []
    page_title = (layout.get("page_title") or "").strip()
    subtitle = (layout.get("subtitle") or "").strip()
    footer_text = (layout.get("footer_text") or "").strip()

    if page_title:
        overlays.append({"type": "title", "text": page_title, "bbox": [0.04, 0.02, 0.70, 0.08]})
    if subtitle:
        overlays.append({"type": "subtitle", "text": subtitle, "bbox": [0.04, 0.095, 0.75, 0.04]})
    if footer_text:
        overlays.append({"type": "footer", "text": footer_text, "bbox": [0.03, 0.92, 0.94, 0.055]})

    return overlays


def build_hybrid_ppt(image_path: str, layout: Dict, output_path: str):
    image_path = str(image_path)
    output_path = str(output_path)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    with Image.open(image_path) as im:
        img_w, img_h = im.size

    margin = Inches(0.10)
    box_left = int(margin)
    box_top = int(margin)
    box_w = int(prs.slide_width - margin * 2)
    box_h = int(prs.slide_height - margin * 2)

    pic_left, pic_top, pic_w, pic_h = _fit_contain(box_left, box_top, box_w, box_h, img_w, img_h)

    slide.shapes.add_picture(image_path, pic_left, pic_top, width=pic_w, height=pic_h)

    overlays: List[Dict] = list(layout.get("overlays") or [])
    if not overlays:
        overlays = _default_overlays(layout)

    for item in overlays[:12]:
        text = (item.get("text") or "").strip()
        kind = item.get("type") or "card_title"
        bbox = item.get("bbox") or []

        if not text or len(bbox) != 4:
            continue

        x, y, w, h = bbox
        if not all(isinstance(v, (int, float)) for v in [x, y, w, h]):
            continue

        # bbox 相对整张原图，映射到 PPT 中图片实际落点
        ox = int(pic_left + x * pic_w)
        oy = int(pic_top + y * pic_h)
        ow = int(w * pic_w)
        oh = int(h * pic_h)

        # 给最小高度，避免太薄
        if oh < Pt(20):
            oh = int(Pt(20))
        if ow < Pt(50):
            ow = int(Pt(50))

        _add_overlay(slide, ox, oy, ow, oh, text, kind)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)
